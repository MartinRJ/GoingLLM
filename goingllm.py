from bs4 import BeautifulSoup
import chardet
import concurrent.futures
from datetime import datetime
from flask import Flask, request, make_response, send_from_directory
import gc
from googleapiclient.discovery import build
from io import BytesIO
import json
from langdetect import detect
import markdown
import mimetypes
from multiprocessing import Process, Manager, Value, Lock
from num2words import num2words
import openai
import openpyxl
import os
import pandas as pd
from pdfminer.high_level import extract_text_to_fp
from pdfminer.layout import LAParams
from pptx import Presentation
import queue
import re
import requests
import spacy
import string
import threading
import tiktoken
import time
from urlextract import URLExtract
from uuid import uuid4
app = Flask(__name__)

from urlextract import URLExtract
import re

SECRET_KEY = os.getenv('SECRETKEY')
CUSTOMSEARCH_KEY = os.getenv('CUSTOMSEARCHKEY')
CX = os.getenv('cx')
TEMPERATURE_DECISION_TO_GOOGLE = float(os.getenv('temperature_decision_to_google'))
MAX_TOKENS_DECISION_TO_GOOGLE = int(os.getenv('max_tokens_decision_to_google'))
TEMPERATURE_CREATE_SEARCHTERMS = float(os.getenv('temperature_create_searchterms'))
MAX_TOKENS_CREATE_SEARCHTERMS = int(os.getenv('max_tokens_create_searchterms'))
TEMPERATURE_SUMMARIZE_RESULT = float(os.getenv('temperature_summarize_result'))
MAX_TOKENS_SUMMARIZE_RESULT = int(os.getenv('SUMMARIZE_MAX_TOKEN_LENGTH'))
MIN_TOKENS_SUMMARIZE_RESULT = int(os.getenv('SUMMARIZE_MIN_TOKEN_LENGTH'))
MAX_FILE_CONTENT = int(os.getenv('MAX_FILE_CONTENT'))
NUMBER_GOOGLE_RESULTS = int(os.getenv('NUMBER_GOOGLE_RESULTS'))
NUMBER_OF_KEYWORDS = int(os.getenv('NUMBER_OF_KEYWORDS'))
TEMPERATURE_FINAL_RESULT = float(os.getenv('temperature_final_result'))
MAX_TOKENS_FINAL_RESULT = int(os.getenv('FINALRESULT_MAX_TOKEN_LENGTH'))
TEMPERATURE_SELECT_SEARCHES = float(os.getenv('temperature_select_searches'))
MAX_TOKENS_SELECT_SEARCHES_LENGTH = int(os.getenv('SELECT_SEARCHES_MAX_TOKEN_LENGTH'))
TEMPERATURE_MORE_SEARCHES = float(os.getenv('temperature_more_searches'))
MAX_TOKENS_MORE_SEARCHES_LENGTH = int(os.getenv('MORE_SEARCHES_MAX_TOKEN_LENGTH'))
TEMPERATURE_CLEANUP = float(os.getenv('temperature_cleanup'))
MAX_TOKENS_CLEANUP_LENGTH = int(os.getenv('CLEANUP_MAX_TOKEN_LENGTH'))
BODY_MAX_LENGTH = int(os.getenv('BODY_MAX_LENGTH'))
GLOBAL_CHATCOMPLETION_TIMEOUT = int(os.getenv('GLOBAL_CHATCOMPLETION_TIMEOUT'))

#SUCCESS/ERROR CODES
FINAL_RESULT_CODE_ERROR_INPUT = "-700" # Error with input
FINAL_RESULT_CODE_ERROR_CHATCOMPLETIONS = "-500" # Error in ChatCompletions API
FINAL_RESULT_CODE_ERROR_CUSTOMSEARCH = "-400" # Error in Custom Search API
FINAL_RESULT_CODE_ERROR_OTHER_CUSTOM = "-600" # Other Error - error message in final_result
FINAL_RESULT_CODE_SUCCESS_WITHOUT_CUSTOMSEARCH = "100" # Success (ChatCompletions-only result)
FINAL_RESULT_CODE_SUCCESS_WITH_CUSTOMSEARCH = "200" # Success (successfully used Custom Search API)
FINAL_RESULT_CODE_CONTINUING_CUSTOMSEARCH = "300" # Working (more search required used Custom Search API)

MESSAGE_MORE_SEARCH_REQUIRED = "More searches required..."

# Texts:
summarytext_newline = "\n"
summarytext_entry_start = "• "
summarytext_entry_end_URL_start = " ["
summarytext_entry_URL_end = "]"

SYSTEM_PROMPT_FINAL_QUERY = "I am \"GoingLLM\" or \"Assistant\", your personal assistant for internet research. As input, I receive a summary of a current internal web crawl. You as a user do not know or see this research information from the requests to me, the research is done internally, you will only ever be able to see my answer and your original request (example: >>What time is it?<<). My answer must not contain direct references to the web crawl and the summaries of the results of the same, as the user will not see these. Instead, I should use the information from the web crawl to improve my answer to your query in a factual and precise way, without referring to incomplete sentences or missing information in the summaries. I answer in the language of the original question."
ASSISTANT_FINAL_QUERY = "Um Ihnen eine angemessene Antwort oder Lösung zu geben, benötige ich weitere Informationen. Könnten Sie bitte mehr Details oder Kontext zu Ihrer Frage oder Aufgabe bereitstellen? Das hilft mir, Ihnen besser zu helfen."

#BASIC AUTHENTICATION
AUTH_UNAME = os.getenv('AUTH_UNAME')
AUTH_PASS = os.getenv('AUTH_PASS')

MODEL = os.getenv('model')
MODEL_MAX_TOKEN = int(os.getenv('model_max_token'))

openai.api_key = SECRET_KEY

@app.route("/", methods=['POST'])
def startup():
    authd = doauthorization()
    if authd:
        return authd
    try:
        body = request.get_data(as_text=True)
        if not body:
            raise ValueError('Empty body')
    except Exception as e:
        debuglog("There was an error with the input.")
        return f'Error extracting body: {e}', 400

    if len(body) > BODY_MAX_LENGTH:
        task_id = str(uuid4())
        errormessage = "Input is too long."
        debuglog(errormessage)
        writefile(FINAL_RESULT_CODE_ERROR_INPUT, errormessage, task_id)
        response = make_response('', 200)
        response.headers['task_id'] = task_id
        return response
    else:
        usertask = body
        dogoogleoverride = False
        always_google = request.headers.get('X-Always-Google')
        if always_google and always_google.lower() == 'true':
            dogoogleoverride = True

        #create new JSON output file with status 'started' and send a 200 response, and start the actual tasks.
        task_id = str(uuid4())
        debuglog(f"New task {task_id} started. User prompt: \"{usertask}\"",True)
        threading.Thread(target=response_task, args=(body, task_id, dogoogleoverride)).start()
        writefile("0", False, task_id)
        response = make_response('', 200)
        response.headers['task_id'] = task_id
        return response

def doauthorization():
    auth = request.authorization #Basic authentication
    if not auth or not (auth.username == AUTH_UNAME and auth.password == AUTH_PASS):
        response = make_response('Could not verify your login!', 401)
        response.headers['WWW-Authenticate'] = 'Basic realm="Login Required"'
        return response
    else:
        return False

@app.route('/searches/<filename>')
def download_file(filename):
    authd = doauthorization()
    if authd:
        return authd
    return send_from_directory('searches', filename)

@app.route('/')
def index():
    return app.send_static_file('assistant.html')

@app.route('/<path:path>')
def static_files(path):
    return send_from_directory(os.path.join(app.root_path, 'static'), path)

@app.route('/tmp/log.txt')
def serve_log_file():
    authd = doauthorization()
    if authd:
        return authd
    try:
        return send_from_directory('tmp', 'log.txt')
    except Exception as e:
        return str(e)

def writefile(progress, json_data, task_id):
    data = {}
    if not json_data:
        data = {
            "task_id": task_id,
            "progress": progress
        }
    else:
        data = {
            "task_id": task_id,
            "progress": progress,
            "answer": json_data
        }
    try:
        # create a 'searches' directory if it does not exist
        if not os.path.exists('searches'):
            os.makedirs('searches')
        # set the file path
        file_path = f'searches/{task_id}.json'
        # open file in write mode
        print(f"Writing to /{file_path}", flush=True)
        with open(file_path, 'w') as f:
            # write JSON data to file
            f.write(json.dumps(data))
    except Exception as e:
        print("Could not write file")

def response_task(usertask, task_id, dogoogleoverride):
    # Preprocess user input
    usertask = preprocess_user_input(usertask)

    dogooglesearch = should_perform_google_search(usertask, dogoogleoverride, task_id) #Should the tool do a google search?

    final_result = ""
    final_result_code = ""
    if dogooglesearch == None or not dogooglesearch:
            if dogooglesearch == None:
                debuglog("Chatcompletions error in should_perform_google_search")
                pass
            else:
                debuglog("No Google search necessary. Generating final response without search results.")
                pass
            final_result, final_result_code = generate_final_result_without_search(usertask, task_id, False)
    elif dogooglesearch:
        debuglog("With Google-search, generating keywords.")
        keywords = generate_keywords(usertask, task_id)
        if keywords == None or not keywords:
            if keywords == None:
                debuglog("Chatcompletions error in generate_keywords")
                pass
            else:
                debuglog("No search terms. Generating final response without search results.")
                pass
            final_result, final_result_code = generate_final_result_without_search(usertask, task_id, True)
        elif valid_keywords(keywords):
            debuglog("Keywords are valid, starting search.")
            searchresults = process_keywords_and_search(keywords, usertask, task_id)
            if searchresults == None or not searchresults:
                if searchresults == None:
                    debuglog("Chatcompletions error in process_keywords_and_search")
                    pass
                else:
                    debuglog("Searchresults are empty. Generating final response without search results.")
                    pass
                final_result, final_result_code = generate_final_result_without_search(usertask, task_id, True)
            else:
                debuglog("Got search results, generating final results.")
                #Is the information sufficient?
                moresearches = do_more_searches(searchresults, usertask, task_id)
                valid_moresearches = False
                try:
                    moresearches = json.loads(moresearches)
                    valid_moresearches = validate_more_searchresults_json(moresearches)
                except Exception as e:
                    debuglog(f"Error in validating moresearches: {e}")
                    moresearches = None

                if moresearches and valid_moresearches and not detectNo(moresearches):
                    # More searches
                    writefile(FINAL_RESULT_CODE_CONTINUING_CUSTOMSEARCH, MESSAGE_MORE_SEARCH_REQUIRED, task_id)
                    # First clean up
                    searchresults, moresearches = cleanup(searchresults, usertask, moresearches)
                    more_searchresults = process_more_searchresults_response(moresearches, searchresults, usertask, task_id)
                    if more_searchresults:
                        searchresults.append(more_searchresults)
                final_result = generate_final_response_with_search_results(searchresults, usertask, task_id)
                if final_result == None:
                    debuglog("Chatcompletions error in generate_final_response_with_search_results")
                    final_result_code = FINAL_RESULT_CODE_ERROR_CHATCOMPLETIONS
                else:
                    debuglog("Success with search results")
                    final_result_code = FINAL_RESULT_CODE_SUCCESS_WITH_CUSTOMSEARCH
        else:
            debuglog("Keywords are not valid. Generating final response without search results.")
            final_result, final_result_code = generate_final_result_without_search(usertask, task_id, True)

    #html = markdown.markdown(responsemessage)
    writefile(final_result_code, final_result, task_id)

    gc.collect() #Cleanup

def cleanup(searchresults, usertask, moresearches):
    #Ask GPT which searchresults should be kept, the rest will be removed. Returns the original searchresults on error.
    prompt = f"Bitte räume die folgenden Summaries zur Useranfrage >>{usertask}<< auf. Antworte ausschließlich mit einem JSON-Objekt mit dem Objekt \"cleanedup\", welches die Indexe derjenigen Resultate beinhaltet, die für die finale anschließende Beantwortung der Frage benötigt werden, diejenigen die nicht nötig bzw. hilfreich sind sollen nicht im cleandup-JSON-Objekt aufgelistet werden. Der angegebene Prozent-Wert gibt an wie viel Token anteilsmäßig für die Erstellung der jeweiligen Summary verwendet wurden. \n\nHier sind die bisherigen gesammelten Summaries, die für die Beantwortung der Useranfrage gesammelt wurden:\n\n{json.dumps(searchresults)}"
    system_prompt = "Ich bin ein persönlicher Assistent für die Internet-Recherche. Ich soll die bisher gesammelten Informationen zu einer bestimmten Useranfrage aufräumen. Ich antworte ausschließlich entweder mit einem JSON-Objekt mit dem Objekt \"cleanedup\", welches die Indexe derjenigen Resultate beinhaltet, die für die finale anschließende Beantwortung der Frage benötigt werden. Beispiel: {\"cleanedup\": [\"1\",\"3\",\"4\"]} Ich antworte in jedem Fall ohne weitere Erklärung oder Kommentar. Der User sieht meine Antwort nicht, sie wird nur von einem Skript weiterverarbeitet, um dem User im späteren Programmablauf die zusammengefassten Ergebnisse der Recherchen präsentieren zu können."
    prompt = truncate_string_to_tokens(prompt, MAX_TOKENS_CLEANUP_LENGTH, system_prompt)
    responsemessage = chatcompletion_with_timeout(system_prompt, prompt, TEMPERATURE_CLEANUP, MAX_TOKENS_CLEANUP_LENGTH)
    if not responsemessage:
        return searchresults # (Fatal) error in chatcompletion, return original searchresults
    debuglog(f"Cleanup. Keep only the following searchresults: {responsemessage}")
    keep_json = extract_json_object(responsemessage)
    if not keep_json:
        return searchresults, moresearches
    new_searchresults, new_moresearches = remove_searchresults(searchresults, keep_json, moresearches)
    if not new_searchresults:
        return searchresults, moresearches
    return new_searchresults, new_moresearches

def remove_searchresults(searchresults, keep_json, moresearches):
    # Returns a version of searchresults with every entry removed that is NOT listed in keep_json.
    # Keep only entries that are listed in keep_json.
    # If there is no 'cleanup' object in keep_json, or on error, return False.
    # Update indices in moresearches to reflect the changes in searchresults.
    # Adds links to moresearches, if an index is removed.
    # moresearches must be already validated.

    if keep_json is None:
        debuglog("remove_searchresults - invalid keep_json")
        return False, moresearches
    try:
        if "cleanedup" in keep_json:
            cleanedup_indices = set(map(int, keep_json["cleanedup"]))
            cleaned_searchresults = [searchresults[i] for i in range(len(searchresults)) if i in cleanedup_indices]

            reindexed_searchresults = []
            index_map = {}
            for new_idx, old_idx in enumerate(cleanedup_indices):
                if old_idx < len(searchresults):  # Check if old_idx is within range of searchresults
                    old_key = str(old_idx)
                    result = {key: value for key, value in searchresults[old_idx].items() if key == old_key}
                    reindexed_searchresults.append({str(new_idx): result[old_key]})
                    index_map[old_idx] = new_idx

            # Update the 'documents' list and 'links' list in moresearches
            if "documents" in moresearches:
                updated_documents = []
                for doc_idx in moresearches["documents"]:
                    int_doc_idx = int(doc_idx)
                    if int_doc_idx in index_map:
                        updated_documents.append(str(index_map[int_doc_idx]))
                    else:
                        if int_doc_idx < len(searchresults):  # Check if int_doc_idx is within range of searchresults
                            if "links" not in moresearches:
                                moresearches["links"] = []
                            moresearches["links"].append(searchresults[int_doc_idx][doc_idx]["URL"])
                            if "openLinks" not in moresearches["action"]:
                                moresearches["action"].append("openLinks")
                moresearches["documents"] = updated_documents
            return reindexed_searchresults, moresearches
        else:
            debuglog("remove_searchresults - no \"cleanedup\" object detected")
            return False, moresearches
    except Exception as e:
        debuglog(f"Error in remove_searchresults: {e}")
        debuglog(f"searchresults: {searchresults}, keep_json: {keep_json}, moresearches: {moresearches}")
        return False, moresearches

def check_type(var):
    # For debug purposes, use like this:
    # var_type = check_type(var)
    if isinstance(var, str):
        return "string"
    elif isinstance(var, list):
        return "list"
    elif isinstance(var, dict):
        return "dictionary"
    else:
        return "neither a string, list nor dictionary"

def detectNo(json_obj):
    #Checks whether the answer to: "do you need more research" was No
    try:
        if json_obj is None:
            raise ValueError("No valid JSON object found in the response.")

        if "action" in json_obj:
            if isinstance(json_obj["action"], list):
                if "Nein" in json_obj["action"] or "Nein." in json_obj["action"]:
                    return True
                else:
                    return False
            else:
                raise ValueError("The 'action' key is not a list.")
        else:
            raise ValueError("The 'action' key is missing from the JSON object.")
    except ValueError as e:
        debuglog(f"Error: {e}")
        return False

def extract_json_object(text):
    # Extracts the first JSON object from the given text
    try:
        debuglog(f"extract_json_object - Input text: {text}")
        json_str = re.search(r'\{.*\}', text).group()
        return json.loads(json_str)
    except (AttributeError, json.JSONDecodeError):
        debuglog("Error in extract_json_object")
        return None

def process_more_searchresults_response(json_object, searchresults, usertask, task_id):
    # Calculate remaining tokens
    # Format of the json_object: {"action": ["<Action type 1>", "<Action type 2>"], "keywords": ["<Keyword1>", "<Keyword2>", ...], "documents": ["0", "1", ...], "links": ["<Link1>", "<Link2>", ...]}
    # Action types: "search", "viewDocuments", "openLinks". One action type minimum. Everything else is optional. If "search" is present, the "keywords" key with value(s) must be present. If "viewDocuments" is present, the "documents" key with value(s) must be present. If "openLinks" is present, the "links" key with value(s) must be present.
    try:
        string = ''.join([f'{summarytext_newline}{summarytext_entry_start}{searchresults[i][str(i)]["summary"]}{summarytext_entry_end_URL_start}{searchresults[i][str(i)]["URL"]}{summarytext_entry_URL_end}' for i in range(len(searchresults)) if len(searchresults[i][str(i)]["summary"]) > 0])
    except TypeError as e:
        debuglog(f"Error in process_more_searchresults_response: {e}")
        debuglog(f"searchresults: {searchresults}", flush=True)
        return searchresults
    added = False
    current_tokens = calculate_tokens(usertask, SYSTEM_PROMPT_FINAL_QUERY, ASSISTANT_FINAL_QUERY, string)
    remaining_tokens = MODEL_MAX_TOKEN - current_tokens - MAX_TOKENS_FINAL_RESULT - calculate_tokens(f"{summarytext_newline}{summarytext_entry_start}{summarytext_entry_end_URL_start}{summarytext_entry_URL_end}")

    if remaining_tokens < MIN_TOKENS_SUMMARIZE_RESULT:
        debuglog(f"Not enough tokens left for another searchresult entry.")
        return searchresults

    debuglog(f"More searches, previous (cleaned-up) results are: {searchresults}")
    debuglog(f"Remaining tokens: {str(remaining_tokens)}")
    #Detect actions
    actions = json_object["action"]
    if "search" in actions:
        keywords = json_object["keywords"]
        if keywords:
            # Perform the search action with the given keywords
            debuglog(f"Searching for: {keywords}")
            new_searchresults = process_keywords_and_search(keywords, usertask, task_id)
            if new_searchresults == None or not new_searchresults:
                if new_searchresults == None:
                    debuglog("Chatcompletions error in process_keywords_and_search")
                    pass
                else:
                    debuglog("Searchresults are empty. Generating final response without search results.")
                    pass
            else:
                if isinstance(more_searchresults):
                    more_searchresults.append(new_searchresults)
                else:
                    more_searchresults = new_searchresults
                    added = True
        else:
            debuglog("No keywords specified. Search is not performed.")
            pass

    if "viewDocuments" in actions:
        documents = json_object["documents"]
        if documents:
            # Perform the action to display documents
            debuglog(f"Display documents: {documents}")
            pass
        else:
            debuglog("No documents specified. Displaying documents is not performed.")
            pass

    if "openLinks" in actions:
        links = json_object["links"]
        if links:
            # Perform the action to open links
            debuglog(f"Open links: {links}")
            pass
        else:
            debuglog("No links specified. Opening links is not performed.")
            pass

    if added:
        return more_searchresults
    else:
        return False

def validate_more_searchresults_json(response_json):
    action_types = {"search", "viewDocuments", "openLinks", "Nein"}

    if not isinstance(response_json, dict):
        raise ValueError("The response JSON object must be a dictionary.")

    if "action" not in response_json:
        raise ValueError("The response JSON object is missing the required 'action' key.")

    actions = response_json["action"]
    if not isinstance(actions, list) or not all(isinstance(action, str) for action in actions):
        raise ValueError("The 'action' list must be a list of strings.")

    for action in actions:
        if action not in action_types:
            raise ValueError(f"Invalid action: {action}. Action type must be one of the following: {action_types}")

    if "Nein" in actions and len(actions) > 1:
        raise ValueError("If 'Nein' is present in the 'action' list, it must be the only action.")

    if not any(key in response_json for key in ("keywords", "documents", "links")):
        raise ValueError("The response JSON object must have at least one of the following keys: 'keywords', 'documents', 'links'")

    if "keywords" in response_json:
        keywords = response_json["keywords"]
        if not isinstance(keywords, list) or not all(isinstance(keyword, str) for keyword in keywords):
            raise ValueError("The 'keywords' list must be a list of strings.")

    if "documents" in response_json:
        documents = response_json["documents"]
        if not isinstance(documents, list) or not all(isinstance(document, str) and document.isdigit() for document in documents):
            raise ValueError("The 'documents' list must be a list of string-converted integers.")

    if "links" in response_json:
        links = response_json["links"]
        if not isinstance(links, list) or not all(isinstance(link, str) for link in links):
            raise ValueError("The 'links' list must be a list of strings.")

    return True

def do_more_searches(searchresults, usertask, task_id):
    prompt = f"Basierend auf den folgenden Informationen, möchten Sie noch weitere Keywords googeln, mehr von den bereits durchsuchten Dokumenten sehen oder noch andere Links direkt aufrufen? Antworten Sie ausschließlich mit entweder \"Nein\" oder mit einem JSON-Objekt.  Möchten Sie noch weitere Keywords googeln, mehr von den bereits durchsuchten Dokumenten sehen oder noch andere Links direkt aufrufen, so geben Sie bitte ein JSON-Objekt mit Ihren Anforderungen zurück, das wie folgt formatiert ist: \n{{\"action\": [\"<Aktionstyp1>\", \"<Aktionstyp2>\"], \"keywords\": [\"<Keyword1>\", \"<Keyword2>\", ...], \"documents\": [\"0\", \"1\", ...], \"links\": [\"<Link1>\", \"<Link2>\", ...]}}\n\nAktionstyp kann entweder \"search\", \"viewDocuments\" oder \"openLinks\" sein. Wenn nein, antworten Sie ausschließlich mit \"Nein\", ohne weitere Erläuterung. Antworten Sie ausschließlich mit entweder einem JSON-Objekt wie angegeben oder mit \"Nein\" als Element der \"action\"-Liste im JSON-Objekt (z. B. {{\"action\": [\"Nein\"], \"keywords\": [], \"documents\": [], \"links\": []}}).\n\nOriginale Useranfrage: \"{usertask}\"\n\nHier sind die Summaries und die prozentualen Anteile der Tokens, die für die Erstellung der Zusammenfassungen jeweils verwendet wurden. Der angegebene Prozent-Wert gibt an wie viel Token anteilsmäßig für die Erstellung der jeweiligen Summary verwendet wurden:\n"
    system_prompt = "Ich ein persönlicher Assistent für die Internet-Recherche. Ich soll entscheiden, ob zur Bearbeitung einer bestimmten Useranfrage noch zusätzliche Google-Recherchen oder Web-Crawls nötig sind, oder ob die bisher gesammelten Informationen in Kombination mit den Daten aus meinen internen Datenbanken ausreichen. Ich antworte ausschließlich entweder mit einem JSON-Objekt oder mit \"Nein\" als Element der \"action\"-Liste im JSON-Objekt (z. B. {\"action\": [\"Nein\"], \"keywords\": [], \"documents\": [], \"links\": []}), falls die Informationen (die ich selbst in einer Anfrage zuvor als Summary zusammengefasst habe) ausreichen. Das JSON-Objekt soll folgendes Format haben (Beispiel): \"{\"action\": [\"<meine auswahl - optional>\" oder \"Nein\",\"<meine auswahl - optional>\",\"<meine auswahl - optional>\"], \"keywords\": [\"<optional - keywords>\"], \"documents\": [\"<optional - index des Dokuments>\"], \"links\": [\"<optional - Direktlinks>\"]}\". Ich antworte in jedem Fall ohne weitere Erklärung oder Kommentar. Der User sieht meine Antwort nicht, sie wird nur von einem Skript weiterverarbeitet, um dem User im späteren Programmablauf die zusammengefassten Ergebnisse der Recherchen präsentieren zu können."
    debuglog(f"Searching with searchresults:\n{json.dumps(searchresults)}")
    prompt = truncate_string_to_tokens(f"{prompt}{json.dumps(searchresults)}", MAX_TOKENS_MORE_SEARCHES_LENGTH, system_prompt)
    responsemessage = chatcompletion_with_timeout(system_prompt, prompt, TEMPERATURE_MORE_SEARCHES, MAX_TOKENS_MORE_SEARCHES_LENGTH)
    debuglog(f"Response of more search:\n{responsemessage}")
    if not responsemessage:
        return None # (Fatal) error in chatcompletion
    return responsemessage

def generate_final_result_without_search(usertask, task_id, regular):
    #Perform and evaluate final regular request (without searchresults). 'regular' determines whether this was called due to an error (regular=False).
    final_result = generate_final_response_without_search_results(usertask, task_id, regular)
    if final_result == None:
        debuglog("Chatcompletions error in generate_final_response_without_search_results")
        final_result_code = FINAL_RESULT_CODE_ERROR_CHATCOMPLETIONS
    else:
        debuglog("Success without search results")
        final_result_code = FINAL_RESULT_CODE_SUCCESS_WITHOUT_CUSTOMSEARCH
    return final_result, final_result_code

def generate_final_response_without_search_results(usertask, task_id, regular):
    #Make a regular query
    system_prompt = "Ich bin dein persönlicher Assistent für die Internetrecherche, und antworte gerade ohne Internetrecherche, da ich zuvor entschieden habe, dass die Anfrage keine Internetrecherche benötigt."
    if not regular:
        system_prompt = "Ich bin dein persönlicher Assistent für die Internetrecherche, und muss auf deine Anfrage leider gerade ohne Internetrecherche antworten. Ich hatte zwar zuvor entschieden, dass zur Beantwortung deiner Anfrage eine Internetrecherche nötig oder hilfreich wäre, jedoch gab es leider einen Fehler bei der Internetrecherche."
    usertask = truncate_string_to_tokens(usertask, MAX_TOKENS_FINAL_RESULT, system_prompt)
    final_result = chatcompletion_with_timeout(system_prompt, usertask, TEMPERATURE_FINAL_RESULT, MAX_TOKENS_FINAL_RESULT)
    return final_result

def generate_final_response_with_search_results(searchresults, usertask, task_id):
    try:
        finalquery = ''.join([f'{summarytext_newline}{summarytext_entry_start}{searchresults[i][str(i)]["summary"]}{summarytext_entry_end_URL_start}{searchresults[i][str(i)]["URL"]}{summarytext_entry_URL_end}' for i in range(len(searchresults)) if len(searchresults[i][str(i)]["summary"]) > 0])
    except TypeError as e:
        debuglog(f"Error in generate_final_response_with_search_results: {e}")
        debuglog(f"searchresults: {searchresults}")
        return None

    #debuglog(f"final query - untruncated: finalquery: \"{finalquery}\", system_prompt: \"{SYSTEM_PROMPT_FINAL_QUERY}\"") #----Debug Output
    finalquery = truncate_string_to_tokens(finalquery, MAX_TOKENS_FINAL_RESULT, SYSTEM_PROMPT_FINAL_QUERY)
    finalquery = truncate_at_last_period_or_newline(finalquery) # make sure the last summary also ends with period or newline.
    final_result = chatcompletion_with_timeout(SYSTEM_PROMPT_FINAL_QUERY, usertask, TEMPERATURE_FINAL_RESULT, MAX_TOKENS_FINAL_RESULT, assistantmessage=ASSISTANT_FINAL_QUERY, prompt2=finalquery)
    return final_result

def process_keywords_and_search(keywords, usertask, task_id):
    with Manager() as manager:
        ALLURLS = manager.dict()
        ALLURLS_lock = Lock()
        counter = Value('i', 0)
        counter_lock = Lock()
        searchresults = manager.list()

        processes = []
        for keyword in keywords:
            process = Process(target=customsearch, args=(keyword, usertask, task_id, counter, counter_lock, ALLURLS, ALLURLS_lock, searchresults))
            process.start()
            processes.append(process)
        
        for process in processes:
            process.join()
        # Convert searchresults to a regular list inside the with block
        debuglog(f"Here is the collection of searchresults BEFORE converting into a list and fixing indexes: {searchresults}")
        searchresults_list = list(searchresults)
        searchresults_list = fix_key_order(searchresults_list)
    return searchresults_list

def fix_key_order(searchresults_list):
    # Fix key order of searchresults_list. Because due to the parallel processing, the order of the string-index searchresults_list[i][str(i)] is not guaranteed.
    # Iterate through the list with searchresults_list[i] and put the value of i into each searchresult (searchresults_list[i][str(i)])
    # searchresult has always only one entry (internal convention, for now)
    fixed_searchresults_list = []
    for i in range(len(searchresults_list)):
        searchresult = searchresults_list[i]
        new_searchresult = {str(i): searchresult[next(iter(searchresult))]} # Get the first entry of searchresult and change the key to i
        fixed_searchresults_list.append(new_searchresult)
    return fixed_searchresults_list

def customsearch(keyword, usertask, task_id, counter, counter_lock, ALLURLS, ALLURLS_lock, searchresults):
    search_google_result = search_google(keyword)
    debuglog(f"Search Google result contains the following data: {json.dumps(search_google_result)}") #debug
    if search_google_result is None:
        search_google_result = {"searchresults":[]} # Create a new empty list, if it was empty.

    google_result = []

    # Check for links in the original task
    urls = extract_document_urls_using_urlextract(usertask)
    if len(urls) > 0:
        # use a list comprehension to add https:// to each url if needed
        urls = ["https://" + url if not url.startswith("https://") else url for url in urls]

        for url in urls:
            debuglog(f"URL detected in original request: {url}") #debug
            url_exists = False

            # Add the URL to google_result if it doesn't already exist
            if url not in google_result:
                google_result.insert(0, url)

            # Check if the URL exists in search_google_result
            for search_result in search_google_result['searchresults']:
                for key, value in search_result.items():
                    if value['url'] == url:
                        url_exists = True
                        break

            # If the URL doesn't exist in search_google_result, add it as a new entry
            if not url_exists:
                new_entry = {
                    str(len(search_google_result["searchresults"])): {
                        "title": "Unknown",
                        "url": url,
                        "description": "This URL was provided with the original user prompt."
                    }
                }
                search_google_result["searchresults"].append(new_entry)

    debuglog(f"Search Google result contains the following data after adding manual URLs from the user prompt: {json.dumps(search_google_result)}") #debug
    if len(search_google_result["searchresults"]) < 1: #Skip if nothing was found or there was an error in search
        # The function has returned an error
        debuglog("Nothing was found or there was an error in the search.")
        return
    for search_result in search_google_result['searchresults']:
        for key, value in search_result.items():
            url = value['url']
            debuglog(f"Adding to google_result: {url}")
            # Add the URL to google_result if it doesn't already exist
            if url not in google_result:
                google_result.append(url)

    # Let ChatGPT pick the most promising
    gpturls = False

    prompt = f"Bitte wähle die Reihenfolge der vielversprechendsten Google-Ergebnisse aus der folgenden Liste aus die für dich - würden dir die Inhalte der jeweiligen URL zur Verfügung gestellt - zur Beantwortung der Aufgabe >>{usertask}<< am nützlichsten sein könnten, und gebe sie als JSON-Objekt mit dem Objekt \"weighting\", das index, und einen \"weight\" Wert enthält zurück, der die geschätzte Gewichtung der Relevanz angibt; In Summe soll das den Wert 1 ergeben. Ergebnisse die für die Aufgabe keine Relevanz versprechen, kannst du aus dem resultierenden JSON-Objekt entfernen: \n\n{json.dumps(search_google_result)}\n\nBeispiel-Antwort: {{\"weighting\": {{\"3\":0.6,\"0\":0.2,\"1\":0.1,\"2\":0.1}}}}. Schreibe keine Begründung, sondern antworte nur mit dem JSON-Objekt."
    system_prompt = "Ich bin dein persönlicher Assistent für die Internetrecherche und antworte immer mit JSON-Objekten mit dem Key \"weighting\". Beispiel: {\"weighting\": {\"2\":0.6,\"0\":0.3,\"1\":0.1}}"
    #debuglog(f"Page content - untruncated, prompt: \"{prompt}\", system_prompt: \"{system_prompt}\"") #----Debug Output
    prompt = truncate_string_to_tokens(prompt, MAX_TOKENS_SELECT_SEARCHES_LENGTH, system_prompt)
    #debuglog(f"Page content - truncated, prompt: \"{prompt}\", system_prompt: \"{system_prompt}\"") #----Debug Output
    responsemessage = chatcompletion_with_timeout(system_prompt, prompt, TEMPERATURE_SELECT_SEARCHES, MAX_TOKENS_SELECT_SEARCHES_LENGTH)
    if not responsemessage:
        return None # (Fatal) error in chatcompletion
    weighting = extract_json(responsemessage, "weighting")

    debuglog(f"weighting content: {json.dumps(weighting)}")
    debuglog(f"search_google_result content: {json.dumps(search_google_result)}")
    debuglog(f"google_result content: {json.dumps(google_result)}")
    if weighting:
        # the function returned a dictionary, re-sort
        sorted_weighting = sorted(weighting.items(), key=lambda x: x[1], reverse=True)
        gpturls = {}
        for index, _ in sorted_weighting:
            if int(index) > len(search_google_result['searchresults'])-1:
                break
            gpturls[index] = search_google_result['searchresults'][int(index)][index]['url']
    else:
        # the function returned False, resume unaltered
        debuglog("No results of initial sort.")
        pass #do nothing

    processes = []
    # The function has returned a list of URLs
    for URL in google_result:
        with counter_lock:
            percent = str(counter.value / ((NUMBER_GOOGLE_RESULTS * NUMBER_OF_KEYWORDS)+len(google_result)) * 100)
            counter.value += 1

        writefile(percent, False, task_id)

        # Check if URL is in ALLURLS
        exists_in_queue = False
        with ALLURLS_lock:
            if URL in ALLURLS:
                exists_in_queue = True
            else:
                ALLURLS[URL] = True
        if exists_in_queue:
            continue  # Exists already

        process = Process(target=do_download_and_summary, args=(weighting, gpturls, keyword, usertask, searchresults, URL, task_id))
        process.start()
        processes.append(process)

    for process in processes:
        process.join()

def do_download_and_summary(weighting, gpturls, keyword, usertask, searchresults, URL, task_id):
    #Download file and create a summary, and append to searchresults
    debuglog(f"Here are the URLs: {URL}")
    dlfile = extract_content(URL)
    if not dlfile:
        responsemessage = "Error"
        debuglog(f"Error summarizing URL content: {URL}")
        return
    responsemessage = dlfile

    prompt = (f"Es wurde folgende Anfrage gestellt: >>{usertask}<<. Im Folgenden findest du den Inhalt einer Seite aus den "
            f"Ergebnissen einer Google-Suche zu dieser Anfrage, bitte fasse das Wesentliche zusammen um mit dem Resultat "
            f"die Anfrage später bestmöglich beantworten zu können, stelle sicher, dass du sämtliche relevanten Spezifika, "
            f"die in deinen internen Datenbanken sonst nicht vorhanden sind in der Zusammenfassung erwähnst. Erwähne auch "
            f"die URL oder Webseite wenn sie relevant ist.\n\nVon URL: {URL}\nKeyword: \"{keyword}\"\nInhalt:\n{responsemessage}")
    system_prompt = "Ich bin dein persönlicher Assistent für die Internetrecherche und erstelle präzise Zusammenfassungen von Webseiteninhalten aus Google-Suchergebnissen. Dabei extrahiere ich relevante Informationen und Spezifika, die zur Beantwortung der gestellten Anfrage erforderlich sind und nicht in meinen internen Datenbanken vorhanden sind. Ich erwähne auch die URL oder Webseite, wenn sie relevant ist."

    debuglog(f"Page content - untruncated, prompt: \"{prompt[:200]}\", system_prompt: \"{system_prompt}\"") #----Debug Output

    weighting_value = False
    if gpturls:
        if URL in gpturls.values():
            # Find the corresponding key in the gpturls dictionary
            key = list(gpturls.keys())[list(gpturls.values()).index(URL)]
            # Get the weighting value for the key
            weighting_value = float(weighting[key])

    max_tokens_completion_summarize = MAX_TOKENS_SUMMARIZE_RESULT

    #Check if there's a weighting value for this URL
    if weighting_value and weighting_value > 0 and len(gpturls) > 0:
        calc_tokens = int(MAX_TOKENS_SUMMARIZE_RESULT * len(gpturls) * weighting_value)
        if calc_tokens < MIN_TOKENS_SUMMARIZE_RESULT:
            debuglog(f"Error - not enough tokens left for summary of URL content with weighting: {URL}")
            return;
        max_tokens_completion_summarize = calc_tokens
        debuglog(f"Weighting applied: {str(weighting_value)} weight => {str(max_tokens_completion_summarize)} tokens")
        if max_tokens_completion_summarize < 1:
            max_tokens_completion_summarize = 1 # max_tokens may not be 0

    #Calculate if there are enough tokens left for the current max_tokens_completion_summarize value, otherwise use less.
    #How many tokens are already used up:
    test_finalquery = ''.join([f'{summarytext_newline}{summarytext_entry_start}{searchresults[i][str(i)]["summary"]}{summarytext_entry_end_URL_start}{searchresults[i][str(i)]["URL"]}{summarytext_entry_URL_end}' for i in range(len(searchresults)) if len(searchresults[i][str(i)]["summary"]) > 0])
    sum_results = calculate_tokens(usertask, SYSTEM_PROMPT_FINAL_QUERY, ASSISTANT_FINAL_QUERY, test_finalquery)

    if MODEL_MAX_TOKEN < sum_results + max_tokens_completion_summarize:
        debuglog(f"Decreasing tokens for summary for: {URL}, not enough tokens left: {str(MODEL_MAX_TOKEN - sum_results)}, requested were {str(max_tokens_completion_summarize)}")
        max_tokens_completion_summarize = MODEL_MAX_TOKEN - sum_results #not enough tokens left for the original number of tokens in max_tokens_completion_summarize, use less
        if max_tokens_completion_summarize < MIN_TOKENS_SUMMARIZE_RESULT:
            debuglog(f"Not enough tokens after decreasing, for: {URL}")
            return # Not enough tokens
    if max_tokens_completion_summarize < 1:
        debuglog(f"Error - no tokens left for summary of URL content: {URL}")
        return
    if max_tokens_completion_summarize < MIN_TOKENS_SUMMARIZE_RESULT:
        debuglog(f"Error - not enough tokens left for summary of URL content: {URL}")
        return
    prompt = truncate_string_to_tokens(prompt, max_tokens_completion_summarize, system_prompt)

    responsemessage = chatcompletion_with_timeout(system_prompt, prompt, TEMPERATURE_SUMMARIZE_RESULT, max_tokens_completion_summarize)
    if not responsemessage:
        return None # (Fatal) error in chatcompletion
    responsemessage = truncate_at_last_period_or_newline(responsemessage) #Make sure responsemessage ends with . or newline, otherwise GPT tends to attempt to finish the sentence.
    #debuglog(f"Page content: prompt: \"{prompt}\", system_prompt: \"{system_prompt}\"") #----Debug Output
    #debuglog(f"Page content: result: \"{responsemessage}\", system_prompt: \"{system_prompt}\"") #----Debug Output
    debuglog(f"Appending to searchresults for URL {URL}: {responsemessage}")
    index = len(searchresults)
    weight = weighting_value
    if not weighting_value:
        weight = "unknown"
    else:
        if isinstance(weighting_value, float):
            weight = weighting_value * 100
        else:
            weight = "unknown"

    searchresult = {
        str(index): {
            "URL": URL,
            "summary": responsemessage,
            "prozent": str(weight),
            "keyword": keyword
        }
    }
    searchresults.append(searchresult)
    debuglog(f"Appending successful.")

def valid_keywords(keywords):
    if not keywords:
        debuglog("valid_keywords: (Fatal) error in chat completion")
        return False # (Fatal) error in chatcompletion
    elif all(isinstance(item, str) for item in keywords):
        return True
    else:
        debuglog(f"Not all entries in the keyword-array are strings. Cannot use the results: {json.dumps(keywords)}")
        return False

def generate_keywords(usertask, task_id):
    number_keywords = num2words(NUMBER_OF_KEYWORDS, lang='de')
    number_entries = "einen Eintrag" if NUMBER_OF_KEYWORDS == 1 else f"{number_keywords} Einträge"
    number_searches = "eine Suche" if NUMBER_OF_KEYWORDS == 1 else f"{number_keywords} Suchen"

    prompt = f"Bitte gib das JSON-Objekt als Antwort zurück, das {number_entries} mit dem Schlüssel 'keywords' enthält, mit den am besten geeigneten Suchbegriffen oder -phrasen, um relevante Informationen zu folgender Anfrage mittels einer Google-Suche zu finden: >>{usertask}<<. Wenn die Anfrage dich auffordert nach einer bestimmten Information zu suchen, dann erstelle Suchbegriffe oder -phrasen, welche möglichst genau der Aufforderung in der Anfrage entsprechen. Berücksichtige dabei Synonyme und verwandte Begriffe und ordne die Suchbegriffe in einer Reihenfolge an, die am wahrscheinlichsten zu erfolgreichen Suchergebnissen führt. Berücksichtige, dass die Ergebnisse der {number_searches} in Kombination verwendet werden sollen, also kannst du bei Bedarf nach einzelnen Informationen suchen. Nutze für die Keywords diejenige Sprache die am besten geeignet ist um relevante Suchergebnisse zu erhalten. Für spezifische Suchen verwende Google-Filter wie \"site:\", besonders wenn z.B. nach Inhalten von speziellen Seiten gesucht wird, wie Twitter, in dem Fall suche beispielsweise nach: \"<suchbegriff> site:twitter.com\". Nutze gegebenenfalls auch andere Suchfilter wo immer das helfen kann, zum Beispiel: \"<suchbegriff> filetype:xlsx\", wenn eine Suche nach speziellen Formaten hilfreich ist (hier: Excel-Dateien). Oder wo nötig nutze auch den \"site:\"-Filter um Ergebnisse aus einem bestimmten Land zu finden, zum Beispiel: \"<suchbegriff> site:.de\" um nur Inhalte von Deutschen Seiten zu finden."
    system_prompt = "Ich bin dein persönlicher Assistent für die Internetrecherche, und das Format meiner Antworten ist immer ein JSON-Objekt mit dem Schlüssel 'keywords', das zur Anfrage passende Google-Suchbegriffe oder -phrasen enthält. Ich unterstütze Google-Suchfilter wie site:, filetype:, allintext:, inurl:, link:, related: und cache: sowie Suchoperatoren wie Anführungszeichen, und die Filter after: / before: um Suchergebnisse aus bestimmten Zeiträumen zu finden. Ich berücksichtige besonders spezifische Benutzer-Eingaben in Anfragen. Besonders wenn nach spezifischen Daten oder Formaten verlangt wird, dann passe ich meine auszugebenden Suchbegriffe im JSON-Objekt der Anfrage möglichst genau an. Beispiel-Anwort zu einer Beispiel-Anfrage \"Wie spät ist es?\": {\"keywords\": [\"aktuelle Uhrzeit\",\"Uhrzeit jetzt\",\"Atomuhr genau\"]}."

    prompt = truncate_string_to_tokens(prompt, MAX_TOKENS_CREATE_SEARCHTERMS, system_prompt)
    responsemessage = chatcompletion_with_timeout(system_prompt, prompt, TEMPERATURE_CREATE_SEARCHTERMS, MAX_TOKENS_CREATE_SEARCHTERMS)
    
    if not responsemessage:
        return None # (Fatal) error in chatcompletion

    # Attempt to extract the JSON object from the response
    jsonobject = extract_json(responsemessage, "keywords")
    keywords = jsonobject if jsonobject else [False]

    return keywords

def should_perform_google_search(usertask, dogoogleoverride, task_id):
    #The user can omit the part, where this tool asks Assistant whether it requires a google search for the task
    dogooglesearch = False
    if dogoogleoverride:
        dogooglesearch = True
        return dogooglesearch
    
    # Get current UTC time
    now = datetime.utcnow()
    # Round to the nearest minute
    now = now.replace(second=0, microsecond=0)
    # Format as a string
    now_str = now.strftime("%Y-%m-%d %H:%M")

    prompt = f"Es wurde soeben folgende Anfrage gestellt: >>{usertask}<<. Benötigst du weitere Informationen aus einer Google-Suche, um diese Anfrage im Anschluss zu erfüllen? Bitte antworte mit \"Ja\" oder \"Nein\". Falls du keinen Zugriff auf Informationen hast die notwendig sind um die Anfrage zu beantworten (zum Beispiel falls du nach Dingen wie der aktuellen Uhrzeit oder nach aktuellen Ereignissen gefragt wirst), oder deine internen Informationen in Bezug auf die Anfrage nicht mehr aktuell sind zum aktuellen Zeitpunkt ({now_str} UTC), so antworte mit \"Ja\". Bei Anfragen oder Fragen die du mit dem Wissen aus deinen Datenbanken alleine ausreichend beantworten kannst (zum Beispiel bei der Frage nach der Lösung einfacher Berechnungen wie \"Wieviel ist 2*2?\", die keine zusätzlichen Daten benötigen), antworte mit \"Nein\". Würdest du weitere Recherche-Ergebnisse aus einer Google-Suche benötigen, um diese Anfrage zufriedenstellend zu beantworten, Ja oder Nein?"
    system_prompt = f"Ich bin dein persönlicher Assistent für die Internetrecherche und antworte ausschließlich nur mit \"Ja\" oder \"Nein\" um initial zu entscheiden ob eine zusätzliche Internetsuche nötig sein wird um in Folge eine bestimmte Anfrage zu beantworten. Mir ist bewusst, dass ich zur Lösung der Aufgabe/Anfrage im Verlauf des Chats bei Bedarf mit neuen relevanten Google-Suchresultaten gespeist werde. Für den Fall, dass ich keinen Zugriff auf benötigte Informationen habe die notwendig sind um die Anfrage zu beantworten (zum Beispiel falls nach Dingen wie der aktuellen Uhrzeit oder nach aktuellen Ereignissen gefragt wird), oder meine internen Informationen in Bezug auf eine Anfrage nicht mehr aktuell sind zum aktuellen Zeitpunkt ({now_str} UTC), so antworte ich immer mit \"Ja\", in dem Wissen, dass mir diese Informationen im Verlauf des Chats noch zur Verfügung gestellt werden. Bei Anfragen oder Fragen die ich mit dem Wissen aus meinen Datenbanken alleine ausreichend beantworten kann (zum Beispiel bei der Frage nach der Lösung einfacher Berechnungen wie \"Wieviel ist 2*2?\", die keine zusätzlichen Daten benötigen), antworte ich immer mit \"Nein\"."
    prompt = truncate_string_to_tokens(prompt, MAX_TOKENS_DECISION_TO_GOOGLE, system_prompt)
    responsemessage = chatcompletion_with_timeout(system_prompt, prompt, TEMPERATURE_DECISION_TO_GOOGLE, MAX_TOKENS_DECISION_TO_GOOGLE)
    if not responsemessage:
        return None # (Fatal) error in chatcompletion
    debuglog(f"Does ChatGPT require a Google-Search: {responsemessage}")
    dogooglesearch = yes_or_no(responsemessage)
    return dogooglesearch

# Preprocess user input
def preprocess_user_input(usertask):
    if "<<" in usertask or ">>" in usertask:
        usertask = usertask.replace("<<", "»").replace(">>", "«")
    return usertask

def chatcompletion(system_prompt, prompt, completiontemperature, completionmaxtokens, assistantmessage = None, prompt2 = None):
    # If assistantmessage and prompt2 are present, will create an "ongoing" conversation with prompt as the first usermessage, assistantmessage as the reply, and prompt2 as the second usermessage
    try:
        messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ]
        if assistantmessage and prompt2:
            messages = [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt},
                    {"role": "assistant", "content": assistantmessage},
                    {"role": "user", "content": prompt2}
                ]
            debuglog(f"chatcompletion - assistantmessage:\n{assistantmessage}")
            #debug
            lines = prompt2.split('\n')
            if len(lines) > 1:
                debuglog(f"chatcompletion - prompt2:\n{lines[0]}\n...\n{lines[-1]}")
            else:
                debuglog(f"chatcompletion - prompt2:\n{prompt2}")
        response = openai.ChatCompletion.create(
        model=MODEL,
        temperature=completiontemperature,
        max_tokens=completionmaxtokens,
        messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ]
        )
        debuglog(f"Query completed. Usage = prompt_tokens: {str(response['usage']['prompt_tokens'])}, completion_tokens: {str(response['usage']['completion_tokens'])}, total_tokens: {str(response['usage']['total_tokens'])}\n\nPrompt:\n{prompt[:1000]}")
        return response['choices'][0]['message']['content']
    except Exception as e:
        Errormessage = f"Error occured in chatcompletion: {e}"
        debuglog(Errormessage)
        return False

def chatcompletion_with_timeout(system_prompt, prompt, completiontemperature, completionmaxtokens, timeout=GLOBAL_CHATCOMPLETION_TIMEOUT, assistantmessage = None, prompt2 = None):
    with concurrent.futures.ThreadPoolExecutor() as executor:
        future = executor.submit(chatcompletion, system_prompt, prompt, completiontemperature, completionmaxtokens, assistantmessage, prompt2)

        try:
            response = future.result(timeout)
            return response
        except concurrent.futures.TimeoutError:
            Errormessage = f"Error: chatcompletion timed out after {timeout} seconds"
            debuglog(Errormessage)
            return False
        except Exception as e:
            Errormessage = f"Error occured in chatcompletion_with_timeout: {e}"
            debuglog(Errormessage)
            return False

def truncate_at_last_period_or_newline(text):
    language = detect(text) # Detect the language of the text

    # Create a dictionary of language codes to spacy model names
    spacy_models = {
        "en": "en_core_web_sm",
        "de": "de_core_news_sm",
        "fr": "fr_core_news_sm",
        "es": "es_core_news_sm",
        "pt": "pt_core_news_sm",
        "it": "it_core_news_sm",
        "nl": "nl_core_news_sm",
        "el": "el_core_news_sm"
    }

    # Check if the recognised language has a spacy model
    if language in spacy_models:
        model_loaded = False
        model_name = spacy_models[language] # Get the model name from the dictionary
        try:
            nlp = spacy.load(model_name) # Load the model
        except:
            try:
                debuglog(f"Model not downloaded, downloading {model_name}")
                spacy.cli.download(model_name) # download the models automatically if they are not present
                nlp = spacy.load(model_name) # Load the model
            except:
                debuglog(f"Could not download and load model {model_name}")
                return truncate_legacy(text)

        try:
            doc = nlp(text) # Create a spacy document from the text
            sentences = list(doc.sents) # Create a list of sentences from the document
            last_sentence = sentences[-1] # Get the last sentence from the list
            truncate_index = last_sentence.start_char - 1 # Find the index before the beginning of the last sentence
            return text[:truncate_index] # Cut the text at this index
        except Exception as e:
            debuglog(f"Could not load language. Error: {e}")
            # Fallback to the legacy method
            return truncate_legacy(text)
    else:
        debuglog("This language is not supported by spacy. Using legacy method, truncating at last period or newline.")
        #Use the 'legacy' method, of cutting off at the last period or newline character.
        return truncate_legacy(text)

def truncate_legacy(text):
    last_period = text.rfind('.')
    last_newline = text.rfind('\n')
    # If neither a dot nor an '\n' is found, the text remains unchanged
    if last_period == -1 and last_newline == -1:
        return text
    # Cut off the text at the point or '\n' that occurs later on
    truncate_index = max(last_period, last_newline)
    if truncate_index == last_period:
        # Add 1 to keep the dot in the text
        return text[:truncate_index + 1]
    else:
        return text[:truncate_index] #do not keep the '\n'

def extract_json(stringwithjson, objectname):
    # Find the start and end indices of the outermost JSON object
    start = -1
    end = -1
    brace_count = 0
    for i, char in enumerate(stringwithjson):
        if char == '{':
            if start == -1:
                start = i
            brace_count += 1
        elif char == '}':
            brace_count -= 1
            if brace_count == 0:
                end = i + 1
                break

    #find the JSON object
    if start == -1 or end == 0:
        debuglog("Error: JSON object not found")
        return False

    json_string = stringwithjson[start:end]

    #parse the JSON object
    try:
        # Convert integer keys to strings
        json_string = re.sub(r'\"(\d+)\":', lambda match: '"' + str(match.group(1)) + '":', json_string)
        data = json.loads(json_string)
    except ValueError as e:
        debuglog(f"Error: Malformed JSON object: {stringwithjson}")
        return False

    items = []
    #access the array
    if objectname in data:
        items = data[objectname]
    else:
        debuglog(f"Error: JSON object doesn't contain \"{objectname}\" array: {stringwithjson}")
        return False

    #return the result
    return items

def calculate_available_tokens(token_reserved_for_response):
    #Calculates the available tokens for a request, taking into account the Tokens reserved for the response
    if token_reserved_for_response > MODEL_MAX_TOKEN:
        return 0
    else:
        return MODEL_MAX_TOKEN - token_reserved_for_response

def calculate_tokens(string, system_prompt = None, assistantmessage = None, prompt2 = None):
    # Calculate tokens. Set system_prompt to False to only count a single string, otherwise the entire message will be counted.
    # If assistantmessage and prompt2 are set, the tokens will be calculated for a 2nd prompt with assistant reply in-between.
    try:
        enc = tiktoken.encoding_for_model(MODEL)
    except KeyError:
        enc = tiktoken.get_encoding("cl100k_base")
        debuglog(f"Error using \"{MODEL}\" as encoding model in truncation, falling back to cl100k_base.")

    if system_prompt:
        messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": string}
        ]
        if assistantmessage and prompt2:
            messages.append({"role": "assistant", "content": assistantmessage})
            messages.append({"role": "user", "content": prompt2})

        num_tokens = 0
        for message in messages:
            num_tokens += 4  # every message follows <im_start>{role/name}\n{content}<im_end>\n
            for key, value in message.items():
                num_tokens += len(enc.encode(value))
                if key == "name":  # if there's a name, the role is omitted
                    num_tokens += -1  # role is always required and always 1 token
        num_tokens += 2  # every reply is primed with <im_start>assistant
        return num_tokens
    else:
        return len(enc.encode(string))

def truncate_string_to_tokens(string, max_tokens, system_prompt):
    # Truncate string to specified number of tokens, if required.
    # max_tokens is what is reserved for the completion (max), string is the user message content, and system_prompt is the system message content.
    base_tokens = 13 #Base value, I noticed that the max is off by 12-13 in the gpt-3.5 API
    try:
        enc = tiktoken.encoding_for_model(MODEL)
    except KeyError:
        enc = tiktoken.get_encoding("cl100k_base")
        debuglog(f"Error using \"{MODEL}\" as encoding model in truncation, falling back to cl100k_base.")

    messages = [
    {"role": "system", "content": system_prompt},
    {"role": "user", "content": string}
    ]
    num_tokens = 0
    for message in messages:
        num_tokens += 4  # every message follows <im_start>{role/name}\n{content}<im_end>\n
        for key, value in message.items():
            num_tokens += len(enc.encode(value))
            if key == "name":  # if there's a name, the role is omitted
                num_tokens += -1  # role is always required and always 1 token
    num_tokens += 2  # every reply is primed with <im_start>assistant

    system_tokens = len(enc.encode(system_prompt))
    possible_tokens = MODEL_MAX_TOKEN - max_tokens - system_tokens - base_tokens
    if (num_tokens > possible_tokens):
        debuglog(f"Length: {str(num_tokens)} tokens. Too long, truncating to {str(possible_tokens)}")
        tokens = enc.encode(string)
        truncated_tokens = tokens[:possible_tokens] # truncate the tokens if they exceed the maximum
        truncated_string = enc.decode(truncated_tokens) # decode the truncated tokens
        return truncated_string
    else:
        debuglog(f"Length: {str(num_tokens)} tokens. Resuming.")
        return string

def yes_or_no(string):
  # Define a boolean variable as return value
  ausgabe = True

  # Try to check the beginning of the string
  try:
    if string.startswith("Nein"):
      ausgabe = False
    else:
      # Assume "Yes" by default
      ausgabe = True
  # Catch an error if the string is not a valid parameter
  except AttributeError:
    # Do nothing and return True
    pass

  # Return the output
  return ausgabe

def search_google(query):
    # Initialise the API with your key and search engine
    service = build("customsearch", "v1", developerKey=CUSTOMSEARCH_KEY)
    debuglog(f"Google search for: \"{query}\"")
    cse = service.cse()
    try:
        # Make a search request to the API
        response = cse.list(q=query, cx=CX).execute()

        # Check if there are search results
        if "items" in response:
            # Extract the first three URLs from Google search results or less if there are not enough
            results = {"searchresults":[]}
            count = 0
            for item in response["items"][:min(NUMBER_GOOGLE_RESULTS, len(response["items"]))]:
                result = {
                    str(count): {"title": item["title"],
                    "url": item["link"],
                    "description": item["snippet"]}
                }
                results["searchresults"].append(result)
                count += 1 # increment count for each result
            return results
        else:
            # There were no search results for this query
            debuglog("No search results for this query.")
            return None
    except Exception as e:
        debuglog(f"Error in Google API query: {e}")
        return None

def load_url_text(url):
    headers = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/89.0.4389.82 Safari/537.36"
    }
    try:
        with requests.get(url, headers=headers, timeout=(3, 8), allow_redirects=True) as response:
            response.raise_for_status()
            # process response
            status_code = response.status_code
            if status_code == 200:
                text = response.text
                if len(text) > 0:
                    return text
                else:
                    return False
            else:
                return False
    except requests.exceptions.Timeout:
        debuglog("Request timed out")
        return False
    except requests.exceptions.RequestException as e:
        debuglog(f"Request error ind load_url_text: {e}")
        return False

def load_url_content(url):
    headers = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/89.0.4389.82 Safari/537.36"
    }
    try:
        with requests.get(url, headers=headers, timeout=(3, 8), allow_redirects=True) as response:
            response.raise_for_status()
            # process response
            status_code = response.status_code
            if status_code == 200:
                content = response.content
                if len(content) > 0:
                    return content
                else:
                    return False
            else:
                return False
    except requests.exceptions.Timeout:
        debuglog("Request timed out")
        return False
    except requests.exceptions.RequestException as e:
        debuglog(f"Request error in load_url_content: {e}")
        return False

def replace_newlines(text):
    return re.sub(r'\n{4,}', '\n', text) #replace all occurrences of four newlines

# Define a function that takes a URL as a parameter and extracts the content
def extract_content(url):
    # Try to send a request to the URL and catch possible exceptions
    # Add a list of supported file extensions
    supported_extensions = [".md", ".log", ".conf", ".config", ".ini", ".yml", ".yaml", ".xml", ".json", ".html", ".php", ".js", ".py", ".java", ".c", ".cpp", ".cs", ".rb", ".sh", ".r", ".m", ".sql"]
    url_extension = os.path.splitext(url)[-1]
    mimetype, encoding = mimetypes.guess_type(url)
    headers = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/89.0.4389.82 Safari/537.36"
    }
    try:
        with requests.head(url, headers=headers, timeout=(3, 8), allow_redirects=True) as response:
            response.raise_for_status()
    except requests.exceptions.Timeout:
        debuglog("Request timed out")
        return False
    except requests.exceptions.RequestException as e:
        debuglog(f"Request error in extract_content: {e}")
        return False
    else:
        # process response
        status_code = response.status_code

        try:
            # Check the status code of the response
            if mimetype is None:
                debuglog(f"Could not determine mimetype for URL: {url}")
                mimetype = response.headers.get("content-type")

            if status_code == 200:
                # Check the content type of the response and handle it accordingly
                if "application/pdf" in mimetype:
                    # Process PDF content
                    with requests.get(url, headers=headers, stream=True, allow_redirects=True) as response:
                        response.raise_for_status()
                        with BytesIO() as filecontent:
                            for chunk in response.iter_content(chunk_size=8192):
                                filecontent.write(chunk)
                            filecontent.seek(0)
                            with BytesIO() as outfp:
                                extract_text_to_fp(filecontent, outfp, laparams=LAParams())
                                text = outfp.getvalue().decode('utf-8')
                                text = replace_newlines(text)
                                debuglog(f"downloaded pdf file: {text[:300]}") #debug
                                return text[:MAX_FILE_CONTENT]
                elif "text/html" in mimetype:
                    filecontent = load_url_text(url)
                    if filecontent: 
                        # Process HTML content
                        # Create a BeautifulSoup object from the HTML string
                        soup = BeautifulSoup(filecontent, "html.parser")
                        html = process_html_content(soup)
                        debuglog(f"downloaded html file: {html[:300]}") #debug
                        return html
                    else:
                        return False
                elif "text/plain" in mimetype:
                    filecontent = load_url_text(url)
                    if filecontent:
                        # Process plain text content
                        filecontent = replace_newlines(filecontent)
                        debuglog(f"downloaded plaintext file: {filecontent[:300]}") #debug
                        return filecontent[:MAX_FILE_CONTENT]
                    else:
                        return False
                elif any(substring in mimetype for substring in ["application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel.sheet.macroEnabled.12"]):
                    # Process Excel content
                    filecontent = load_url_content(url)
                    if filecontent:
                        text = process_excel_content(filecontent)
                        if text:
                            debuglog(f"downloaded excel file: {text[:300]}") #debug
                            return text
                        else:
                            return False
                    else:
                        return False
                elif "text/csv" in mimetype:
                    # Process CSV content
                    filecontent = load_url_content(url)
                    if filecontent:
                        text = process_csv_content(filecontent)
                        if text:
                            debuglog(f"downloaded csv file: {text[:300]}") #debug
                            return text
                        else:
                            return False
                    else:
                        return False
                elif any(substring in mimetype for substring in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint.presentation.macroEnabled.12"]):
                    # Process PowerPoint content
                    filecontent = load_url_content(url)
                    if filecontent:
                        text = process_ppt_content(filecontent)
                        if text:
                            debuglog(f"downloaded powerpoint file: {text[:300]}") #debug
                            return text
                        else:
                            return False
                    else:
                        return False
                elif url_extension in supported_extensions:
                    # If the URL's file extension is in the supported list
                    filecontent = load_url_text(url)
                    if filecontent:
                        # Process the file as plain text
                        filecontent = replace_newlines(filecontent)
                        debuglog(f"downloaded supported file as plaintext: {filecontent[:300]}")  # debug
                        return filecontent[:MAX_FILE_CONTENT]
                    else:
                        return False
                else:
                    # The content type is not supported
                    debuglog(f"Content type '{mimetype}' not supported")
                    return False
            else:
                # The URL could not be found or there was another error
                debuglog(f"Error retrieving URL: {status_code}")
                return False
        except Exception as e:
            # There was another error
            debuglog(f"Error retrieving URL: {e}")
            return False

def process_excel_content(filecontent):
    try:
        # Detect the encoding of the file content using chardet
        with BytesIO(filecontent) as f:
            df = pd.read_excel(f)
            text = df.to_string()
            text = replace_newlines(text)
            return text[:MAX_FILE_CONTENT]
    except Exception as e:
        debuglog(f"Error processing Excel content: {e}")
        return False

def process_csv_content(filecontent):
    try:
        # Detect the encoding of the file content using chardet
        detected_encoding = chardet.detect(filecontent)['encoding']

        with BytesIO(filecontent) as f:
            df = pd.read_csv(f, encoding=detected_encoding)
            text = df.to_string()
            text = replace_newlines(text)
            return text[:MAX_FILE_CONTENT]
    except Exception as e:
        debuglog(f"Error processing CSV content: {e}")
        return False

def process_ppt_content(filecontent):
    try:
        with BytesIO(filecontent) as f:
            pr = Presentation(f)
            text_list = []
            for slide in pr.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_list.append(shape.text)
                        if sum(len(s) for s in text_list) > MAX_FILE_CONTENT:
                            break
            text = "\n".join(text_list)
            text = replace_newlines(text)
            return text[:MAX_FILE_CONTENT]
    except Exception as e:
        debuglog(f"Error processing PowerPoint content: {e}")
        return False

def process_html_content(soup):
    # Find the body element in the HTML document
    body = soup.body
    # Extract the text from the body element
    html = body.get_text()
    html = replace_newlines(html)
    return html[:MAX_FILE_CONTENT]

def debuglog(text, create=False):
    try:
        writemode = 'a'
        if create:
            writemode = 'w' #Overwrite
        # create a 'tmp' directory if it does not exist
        if not os.path.exists('tmp'):
            os.makedirs('tmp')
        # set the file path
        file_path = f'tmp/log.txt'
        # open file in write mode
        with open(file_path, writemode) as f:
            # write text to file
            f.writelines([text, "\n--------------------\n"])
    except Exception as e:
        print(f"Error debuglog: could not write to file: {e}", flush=True)

#----- Extract URL functions
def remove_punctuation(urls):
    cleaned_urls = []
    for url in urls:
        cleaned_url = url.rstrip(string.punctuation)
        cleaned_urls.append(cleaned_url)
    return cleaned_urls

def normalize_urls_protocol(urls):
    normalized_urls = []
    for url in urls:
        # Replace improper protocols and "http://" with "https://"
        url = re.sub(r'^(?:htp://|http:/*|http://|https?://)', 'https://', url, flags=re.IGNORECASE)
       
        # Add "https://" protocol if missing
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url

        normalized_urls.append(url)
    return normalized_urls

def find_additional_urls(text):
    tlds = [".com", ".tk", ".cn", ".de", ".net", ".uk", ".org", ".nl", ".ru", ".br", ".au",
    ".fr", ".eu", ".it", ".pl", ".in", ".info", ".es", ".ca", ".io", ".gov", ".gov.uk", ".gouv.fr",
    ".gc.ca", ".gov.au", ".gov.in", ".gov.za", ".gov.cn", ".gov.br", ".gov.sg", ".gov.de",
    ".gov.it", ".gov.nl", ".gov.my", ".gov.ru", ".gov.ua", ".gov.ie", ".gov.nz", ".gov.il",
    ".gov.pl", ".edu", ".ac", ".sci", ".research", ".scholar", ".ir", ".ch", ".at", ".be", ".dk",
    ".fi", ".gr", ".hu", ".no", ".pt", ".se", ".co", ".biz", ".co.uk", ".ac.uk", ".edu.au",
    ".edu.sg", ".edu.de", ".coop", ".museum", ".pro", ".name", ".govt.nz"]
    additional_urls = []

    # Find URLs with https:// or http:// and special characters at the end
    pattern = re.compile(r'https?://[^\s!\"\'\[\]\(\)]+')
    matches = pattern.finditer(text)
    for match in matches:
        additional_urls.append(match.group())

    # Find URLs without https:// or http://
    url_chars = r"[A-Za-z0-9\-\._~:/\?#\[\]@!\$&'\(\)\*\+,;=äöüÄÖÜß]+"
    
    for tld in tlds:
        pattern = re.compile(fr"{url_chars}{re.escape(tld)}")
        matches = pattern.finditer(text)
        for match in matches:
            if not match.group().startswith(("http://", "https://")):
                additional_urls.append(match.group())
    return additional_urls

def remove_extracted_urls_from_text(text, extracted_urls):
    cleaned_text = text
    for url in extracted_urls:
        index = text.find(url)
        if index > 0 and text[index - 1] == ".":
            cleaned_text = cleaned_text.replace("." + url, "", 1)
        else:
            cleaned_text = cleaned_text.replace(url, "", 1)
    return cleaned_text

def extract_document_urls_using_urlextract(text):
    if not text:
        return []
    #Extracts all URLs from text
    tlds = [".com", ".tk", ".cn", ".de", ".net", ".uk", ".org", ".nl", ".ru", ".br", ".au",
    ".fr", ".eu", ".it", ".pl", ".in", ".info", ".es", ".ca", ".io", ".gov", ".gov.uk", ".gouv.fr",
    ".gc.ca", ".gov.au", ".gov.in", ".gov.za", ".gov.cn", ".gov.br", ".gov.sg", ".gov.de",
    ".gov.it", ".gov.nl", ".gov.my", ".gov.ru", ".gov.ua", ".gov.ie", ".gov.nz", ".gov.il",
    ".gov.pl", ".edu", ".ac", ".sci", ".research", ".scholar", ".ir", ".ch", ".at", ".be", ".dk",
    ".fi", ".gr", ".hu", ".no", ".pt", ".se", ".co", ".biz", ".co.uk", ".ac.uk", ".edu.au",
    ".edu.sg", ".edu.de", ".coop", ".museum", ".pro", ".name", ".govt.nz"]
    file_exts = [
        ".txt", ".csv", ".md", ".log", ".conf", ".config", ".ini", ".yml", ".yaml",
        ".xml", ".json", ".html", ".php", ".js", ".py", ".java", ".c", ".cpp",
        ".cs", ".rb", ".sh", ".r", ".m", ".sql"
    ]

    replacements = []
    uuids = []

    pattern = re.compile(r'/(?P<filename>[\w\-]+)(?P<file_ext>\.[A-Za-z0-9]+)(?P<next_char>[\s\.,:;\?!])')
    #Make sure that TLDs don't override file extensions
    for tld in tlds:
        for file_ext in file_exts:
            if file_ext.startswith(tld):
                matches = pattern.finditer(text)
                for match in matches:
                    if match.group('file_ext') == file_ext:
                        tempuuid = uuid4()
                        temp_tld = f"{tempuuid}.com"
                        original_filename = match.group('filename')
                        replacements.append((original_filename, file_ext, temp_tld))
                        uuids.append(tempuuid)
                        original_str = f"/{match.group('filename')}{file_ext}{match.group('next_char')}"
                        replaced_str = f"/{tempuuid}.com{match.group('next_char')}"
                        text = text.replace(original_str, replaced_str)

    extractor = URLExtract()
    all_urls = extractor.find_urls(text)
    
    # Filter out the URLs that are just TLDs
    cleaned_urls = []
    for url in all_urls:
        index = text.find(url)
        if index == 0 or text[index - 1] != ".":
            cleaned_urls.append(url)

    # Remove the extracted URLs from the text
    cleaned_text = remove_extracted_urls_from_text(text, all_urls)
    
    # Find additional URLs with special characters at the end
    additional_urls = find_additional_urls(cleaned_text)

    cleaned_urls.extend(additional_urls)
    
    cleaned_urls = normalize_urls_protocol(cleaned_urls)
    cleaned_urls = remove_punctuation(cleaned_urls)

    # Replace temporary TLDs and filenames back to original values
    for original_filename, file_ext, temp_tld in replacements:
        cleaned_urls = [url.replace(temp_tld, f"{original_filename}{file_ext}") for url in cleaned_urls]

    # Remove URLs containing any of the stored UUIDs
    cleaned_urls = [url for url in cleaned_urls if not any(str(tempuuid) in url for tempuuid in uuids)]
    # Remove duplicates
    cleaned_urls = list(set(cleaned_urls))
    return cleaned_urls

if __name__ == "__main__":
    app.run()